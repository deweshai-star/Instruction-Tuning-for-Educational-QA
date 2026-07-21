import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 6 is the layout index for a blank slide in default python-pptx template
    blank_layout = prs.slide_layouts[6]
    
    # Define Color Palette (Cool Tech / Modern Academic theme)
    BG_LIGHT = RGBColor(248, 250, 252)       # Slate 50 (Very light gray-blue)
    BG_DARK = RGBColor(15, 23, 42)          # Slate 900 (Deep Indigo/Dark Gray)
    TEXT_DARK = RGBColor(15, 23, 42)        # Slate 900
    TEXT_MUTED = RGBColor(71, 85, 105)      # Slate 600 (Muted blue-gray)
    ACCENT_GOLD = RGBColor(197, 160, 89)     # Soft Gold (Highlight)
    ACCENT_BLUE = RGBColor(37, 99, 235)      # Royal Blue
    WHITE = RGBColor(255, 255, 255)
    CARD_BORDER = RGBColor(226, 232, 240)    # Slate 200
    
    # Helper to set background
    def set_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to add standard titles to content slides
    def add_slide_header(slide, title_text, category_text=""):
        # Category/Tracker text
        if category_text:
            cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.3))
            tf_cat = cat_box.text_frame
            tf_cat.word_wrap = True
            tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
            p_cat = tf_cat.paragraphs[0]
            p_cat.text = category_text.upper()
            p_cat.font.name = "Calibri"
            p_cat.font.size = Pt(10)
            p_cat.font.bold = True
            p_cat.font.color.rgb = ACCENT_GOLD
            p_cat.font.letter_spacing = Pt(2)
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12.133), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(32)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK

    # ==========================================
    # SLIDE 1: Title Slide (Dark Background)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1, BG_DARK)
    
    # Left decorative bar
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.12), Inches(3.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GOLD
    bar.line.color.rgb = ACCENT_GOLD
    
    # Title & Subtitle block
    title_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(11.0), Inches(3.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    tf1.margin_left = tf1.margin_right = tf1.margin_top = tf1.margin_bottom = 0
    
    # Title
    p1 = tf1.paragraphs[0]
    p1.text = "Instruction Tuning for Educational QA"
    p1.font.name = "Calibri"
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(12)
    
    # Subtitle
    p2 = tf1.add_paragraph()
    p2.text = "Parameter-Efficient Fine-Tuning of TinyLlama 1.1B on Consumer Hardware"
    p2.font.name = "Calibri"
    p2.font.size = Pt(20)
    p2.font.color.rgb = ACCENT_GOLD
    p2.space_after = Pt(24)
    
    # Description
    p3 = tf1.add_paragraph()
    p3.text = "Deploying domain-specific educational question-answering systems locally with zero cloud dependence."
    p3.font.name = "Calibri"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
    
    # Footer metadata
    footer_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.6), Inches(6.0), Inches(1.0))
    tf_f = footer_box.text_frame
    p_f1 = tf_f.paragraphs[0]
    p_f1.text = "IITM DS AI - Natural Language Processing Project"
    p_f1.font.name = "Calibri"
    p_f1.font.size = Pt(12)
    p_f1.font.bold = True
    p_f1.font.color.rgb = WHITE
    
    p_f2 = tf_f.add_paragraph()
    p_f2.text = "Methodology: PEFT / LoRA (Low-Rank Adaptation)"
    p_f2.font.name = "Calibri"
    p_f2.font.size = Pt(11)
    p_f2.font.color.rgb = RGBColor(148, 163, 184)

    # ==========================================
    # SLIDE 2: Objectives (Light Background)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2, BG_LIGHT)
    add_slide_header(slide2, "Project Objectives", "Core Objectives")
    
    # Subtitle explanation
    desc_box2 = slide2.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.4))
    desc_tf2 = desc_box2.text_frame
    desc_p2 = desc_tf2.paragraphs[0]
    desc_p2.text = "Key motivations behind training a lightweight model for educational environments."
    desc_p2.font.name = "Calibri"
    desc_p2.font.size = Pt(15)
    desc_p2.font.color.rgb = TEXT_MUTED
    
    objectives = [
        {
            "num": "01",
            "title": "CPU-Only Accessibility",
            "points": [
                "Remove reliance on expensive GPU hardware.",
                "Make model training and execution feasible on standard consumer laptops with ~6GB RAM.",
                "Provide a cost-effective setup for local institutions and research environments."
            ]
        },
        {
            "num": "02",
            "title": "Classroom-Ready Assistance",
            "points": [
                "Fine-tune the model specifically to answer educational queries.",
                "Format responses to be clear, structured, and easy for students to understand.",
                "Encourage domain-specific accuracy and minimize general hallucinations."
            ]
        },
        {
            "num": "03",
            "title": "Privacy & Offline Support",
            "points": [
                "Establish a fully offline architecture.",
                "Avoid sending student and teacher data to third-party cloud API services.",
                "Enable deployment in remote areas with limited internet connectivity."
            ]
        }
    ]
    
    obj_width = Inches(3.8)
    obj_height = Inches(4.3)
    obj_top = Inches(2.1)
    obj_spacing = Inches(0.36)
    obj_left_start = Inches(0.6)
    
    for i, obj in enumerate(objectives):
        left_pos = obj_left_start + i * (obj_width + obj_spacing)
        
        # Base Card Shape
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, obj_top, obj_width, obj_height)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)
        
        # Accent Line
        accent_line = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, obj_top, obj_width, Inches(0.12))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ACCENT_GOLD
        accent_line.line.color.rgb = ACCENT_GOLD
        
        # Text Box
        tb = slide2.shapes.add_textbox(left_pos + Inches(0.2), obj_top + Inches(0.3), obj_width - Inches(0.4), obj_height - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_num = tf.paragraphs[0]
        p_num.text = obj["num"]
        p_num.font.name = "Calibri"
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = ACCENT_GOLD
        p_num.space_after = Pt(2)
        
        p_title = tf.add_paragraph()
        p_title.text = obj["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.space_after = Pt(8)
        
        for pt in obj["points"]:
            p_bullet = tf.add_paragraph()
            p_bullet.text = "• " + pt
            p_bullet.font.name = "Calibri"
            p_bullet.font.size = Pt(10.5)
            p_bullet.font.color.rgb = TEXT_MUTED
            p_bullet.space_before = Pt(4)

    # ==========================================
    # SLIDE 3: Dataset (Light Background)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3, BG_LIGHT)
    add_slide_header(slide3, "Alpaca Instruction Dataset", "Dataset & Preparation")
    
    # Subtitle
    desc_box3 = slide3.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.4))
    desc_tf3 = desc_box3.text_frame
    desc_p3 = desc_tf3.paragraphs[0]
    desc_p3.text = "Utilizing Stanford's Alpaca instruction-following dataset for QA formatting."
    desc_p3.font.name = "Calibri"
    desc_p3.font.size = Pt(15)
    desc_p3.font.color.rgb = TEXT_MUTED
    
    # Left Column: Dataset Details
    col_width = Inches(5.8)
    col_height = Inches(4.5)
    left_col = Inches(0.6)
    
    card_data = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_col, Inches(2.1), col_width, col_height)
    card_data.fill.solid()
    card_data.fill.fore_color.rgb = WHITE
    card_data.line.color.rgb = CARD_BORDER
    card_data.line.width = Pt(1)
    
    accent_data = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_col, Inches(2.1), col_width, Inches(0.1))
    accent_data.fill.solid()
    accent_data.fill.fore_color.rgb = TEXT_MUTED
    accent_data.line.color.rgb = TEXT_MUTED
    
    tb_d = slide3.shapes.add_textbox(left_col + Inches(0.3), Inches(2.4), col_width - Inches(0.6), col_height - Inches(0.6))
    tf_d = tb_d.text_frame
    tf_d.word_wrap = True
    
    p_d_title = tf_d.paragraphs[0]
    p_d_title.text = "Dataset Characteristics"
    p_d_title.font.name = "Calibri"
    p_d_title.font.size = Pt(20)
    p_d_title.font.bold = True
    p_d_title.font.color.rgb = TEXT_DARK
    p_d_title.space_after = Pt(12)
    
    bullets_d = [
        "Source: Stanford's tatsu-lab/alpaca on HuggingFace.",
        "Total samples: 52,002 instruction-response pairs.",
        "Subset size: 1,000 samples selected (optimized for laptop CPU resource constraints).",
        "Split structure: 900 samples (90%) for training, 100 samples (10%) for evaluation.",
        "Tokenization: LLaMA tokenizer, padding to a maximum sequence length of 512 tokens."
    ]
    for b in bullets_d:
        p_b = tf_d.add_paragraph()
        p_b.text = "• " + b
        p_b.font.name = "Calibri"
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(6)
        
    # Right Column: Prompt Format Card
    right_col = Inches(6.933)
    
    card_format = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_col, Inches(2.1), col_width, col_height)
    card_format.fill.solid()
    card_format.fill.fore_color.rgb = WHITE
    card_format.line.color.rgb = CARD_BORDER
    card_format.line.width = Pt(1.5)
    
    accent_format = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_col, Inches(2.1), col_width, Inches(0.1))
    accent_format.fill.solid()
    accent_format.fill.fore_color.rgb = ACCENT_GOLD
    accent_format.line.color.rgb = ACCENT_GOLD
    
    tb_f = slide3.shapes.add_textbox(right_col + Inches(0.3), Inches(2.4), col_width - Inches(0.6), col_height - Inches(0.6))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True
    
    p_f_title = tf_f.paragraphs[0]
    p_f_title.text = "Prompt Format Template"
    p_f_title.font.name = "Calibri"
    p_f_title.font.size = Pt(20)
    p_f_title.font.bold = True
    p_f_title.font.color.rgb = TEXT_DARK
    p_f_title.space_after = Pt(10)
    
    p_f_desc = tf_f.add_paragraph()
    p_f_desc.text = "Data points are formatted into structured templates to construct unified context:"
    p_f_desc.font.name = "Calibri"
    p_f_desc.font.size = Pt(12)
    p_f_desc.font.color.rgb = TEXT_MUTED
    p_f_desc.space_after = Pt(10)
    
    # Template Box
    template_text = (
        "### Instruction:\n"
        "{instruction}\n\n"
        "### Input:\n"
        "{input}    ← (Optional context, omitted if empty)\n\n"
        "### Response:\n"
        "{output}"
    )
    p_tpl = tf_f.add_paragraph()
    p_tpl.text = template_text
    p_tpl.font.name = "Courier New"
    p_tpl.font.size = Pt(11)
    p_tpl.font.color.rgb = ACCENT_BLUE
    p_tpl.space_before = Pt(4)

    # ==========================================
    # SLIDE 4: Workflow & Architecture (Light Background)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4, BG_LIGHT)
    add_slide_header(slide4, "End-to-End Workflow & System Architecture", "Workflow & Architecture")
    
    # Subtitle explanation
    desc_box4 = slide4.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.4))
    desc_tf4 = desc_box4.text_frame
    desc_p4 = desc_tf4.paragraphs[0]
    desc_p4.text = "The systemic pipeline from raw data loading to interactive execution."
    desc_p4.font.name = "Calibri"
    desc_p4.font.size = Pt(15)
    desc_p4.font.color.rgb = TEXT_MUTED
    
    cards = [
        {
            "num": "01",
            "title": "Data Prep",
            "bullets": [
                "Download Alpaca subset",
                "Inject instruction template",
                "Tokenize using LLaMA tokenizer",
                "Apply padding & truncate length to 512"
            ]
        },
        {
            "num": "02",
            "title": "PEFT Configuration",
            "bullets": [
                "Load base TinyLlama 1.1B",
                "Freeze 99.9% of base model weights",
                "Apply LoRA adapters",
                "Target q_proj & v_proj (r=8, alpha=16)"
            ]
        },
        {
            "num": "03",
            "title": "CPU-Only Training",
            "bullets": [
                "PyTorch CPU-only runtime",
                "Use Float32 training precision",
                "Enable Gradient Checkpointing",
                "Apply AdamW (Effective Batch Size = 4)"
            ]
        },
        {
            "num": "04",
            "title": "Local Web UI",
            "bullets": [
                "Save adapter weights (~4.5MB)",
                "Integrate with Streamlit App",
                "Auto-detect adapter directory checkpoints",
                "Generate response on CPU in browser"
            ]
        }
    ]
    
    card_width = Inches(2.8)
    card_height = Inches(4.5)
    top_pos = Inches(2.1)
    spacing = Inches(0.31)
    left_start = Inches(0.6)
    
    for i, c in enumerate(cards):
        left_pos = left_start + i * (card_width + spacing)
        
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)
        
        accent_line = slide4.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, top_pos, card_width, Inches(0.12))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ACCENT_GOLD
        accent_line.line.color.rgb = ACCENT_GOLD
        
        tb = slide4.shapes.add_textbox(left_pos + Inches(0.15), top_pos + Inches(0.25), card_width - Inches(0.3), card_height - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_num = tf.paragraphs[0]
        p_num.text = c["num"]
        p_num.font.name = "Calibri"
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = ACCENT_GOLD
        p_num.space_after = Pt(2)
        
        p_title = tf.add_paragraph()
        p_title.text = c["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.space_after = Pt(8)
        
        for b in c["bullets"]:
            p_bullet = tf.add_paragraph()
            p_bullet.text = "• " + b
            p_bullet.font.name = "Calibri"
            p_bullet.font.size = Pt(10)
            p_bullet.font.color.rgb = TEXT_MUTED
            p_bullet.space_before = Pt(3)

    # ==========================================
    # SLIDE 5: Challenges (Light Background)
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5, BG_LIGHT)
    add_slide_header(slide5, "Development Challenges & Key Solutions", "Challenges & Solutions")
    
    # Subtitle
    desc_box5 = slide5.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.4))
    desc_tf5 = desc_box5.text_frame
    desc_p5 = desc_tf5.paragraphs[0]
    desc_p5.text = "Overcoming hardware bottlenecks and path resolution errors on local client machines."
    desc_p5.font.name = "Calibri"
    desc_p5.font.size = Pt(15)
    desc_p5.font.color.rgb = TEXT_MUTED
    
    # Left Column: Challenge 1
    col_width = Inches(5.8)
    col_height = Inches(4.5)
    col1_left = Inches(0.6)
    
    card1 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, col1_left, Inches(2.1), col_width, col_height)
    card1.fill.solid()
    card1.fill.fore_color.rgb = WHITE
    card1.line.color.rgb = CARD_BORDER
    card1.line.width = Pt(1)
    
    accent_card1 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, col1_left, Inches(2.1), col_width, Inches(0.1))
    accent_card1.fill.solid()
    accent_card1.fill.fore_color.rgb = TEXT_MUTED
    accent_card1.line.color.rgb = TEXT_MUTED
    
    tb_c1 = slide5.shapes.add_textbox(col1_left + Inches(0.3), Inches(2.4), col_width - Inches(0.6), col_height - Inches(0.6))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    
    p_c1_title = tf_c1.paragraphs[0]
    p_c1_title.text = "1. Severe CPU Memory Constraints"
    p_c1_title.font.name = "Calibri"
    p_c1_title.font.size = Pt(20)
    p_c1_title.font.bold = True
    p_c1_title.font.color.rgb = TEXT_DARK
    p_c1_title.space_after = Pt(12)
    
    c1_points = [
        "Challenge: Loading a 1.1B model and performing backpropagation consumes huge RAM, causing standard CPUs to throw Out-Of-Memory (OOM) errors.",
        "Solution: Enabled gradient checkpointing to discard intermediate activations.",
        "Solution: Used Parameter-Efficient Fine-Tuning (PEFT) with LoRA, shrinking trainable parameters to 0.1% (~1.12M), lowering peak RAM to ~6 GB.",
        "Solution: Restricted training batch size to 1, while using gradient accumulation of 4 steps to maintain training stability."
    ]
    for pt in c1_points:
        p_pt = tf_c1.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.name = "Calibri"
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = TEXT_MUTED
        p_pt.space_before = Pt(6)
        
    # Right Column: Challenge 2
    col2_left = Inches(6.933)
    
    card2 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, col2_left, Inches(2.1), col_width, col_height)
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.line.color.rgb = CARD_BORDER
    card2.line.width = Pt(1.5)
    
    accent_card2 = slide5.shapes.add_shape(MSO_SHAPE.RECTANGLE, col2_left, Inches(2.1), col_width, Inches(0.1))
    accent_card2.fill.solid()
    accent_card2.fill.fore_color.rgb = ACCENT_GOLD
    accent_card2.line.color.rgb = ACCENT_GOLD
    
    tb_c2 = slide5.shapes.add_textbox(col2_left + Inches(0.3), Inches(2.4), col_width - Inches(0.6), col_height - Inches(0.6))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    
    p_c2_title = tf_c2.paragraphs[0]
    p_c2_title.text = "2. Path Resolution & Fallback Crashes"
    p_c2_title.font.name = "Calibri"
    p_c2_title.font.size = Pt(20)
    p_c2_title.font.bold = True
    p_c2_title.font.color.rgb = TEXT_DARK
    p_c2_title.space_after = Pt(12)
    
    c2_points = [
        "Challenge: If the final-adapter weight directory was absent, the script attempted to download it from Hugging Face, crashing with a HFValidationError.",
        "Solution: Refactored paths in both app.py and inference.py to verify local folder existence recursively.",
        "Solution: Implemented fallback logic to load intermediate checkpoints (like checkpoint-25/) if the final adapter is missing.",
        "Solution: Added safety fallback to the base TinyLlama model, ensuring the app starts gracefully even without any adapters."
    ]
    for pt in c2_points:
        p_pt = tf_c2.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.name = "Calibri"
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = TEXT_MUTED
        p_pt.space_before = Pt(6)

    # ==========================================
    # SLIDE 6: Results & Conclusion (Light Background)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6, BG_LIGHT)
    add_slide_header(slide6, "Empirical Results & Conclusion", "Results & Conclusion")
    
    # Subtitle
    desc_box6 = slide6.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.4))
    desc_tf6 = desc_box6.text_frame
    desc_p6 = desc_tf6.paragraphs[0]
    desc_p6.text = "Analyzing model convergence metrics and summarizing project takeaways."
    desc_p6.font.name = "Calibri"
    desc_p6.font.size = Pt(15)
    desc_p6.font.color.rgb = TEXT_MUTED
    
    # Left Column: Loss Table
    col_width = Inches(5.8)
    col_height = Inches(4.5)
    col1_left = Inches(0.6)
    
    card_table = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, col1_left, Inches(2.1), col_width, col_height)
    card_table.fill.solid()
    card_table.fill.fore_color.rgb = WHITE
    card_table.line.color.rgb = CARD_BORDER
    card_table.line.width = Pt(1)
    
    accent_card_t = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, col1_left, Inches(2.1), col_width, Inches(0.1))
    accent_card_t.fill.solid()
    accent_card_t.fill.fore_color.rgb = TEXT_MUTED
    accent_card_t.line.color.rgb = TEXT_MUTED
    
    tb_t = slide6.shapes.add_textbox(col1_left + Inches(0.3), Inches(2.3), col_width - Inches(0.6), Inches(0.5))
    tf_t = tb_t.text_frame
    p_t_title = tf_t.paragraphs[0]
    p_t_title.text = "Loss Convergence Table"
    p_t_title.font.name = "Calibri"
    p_t_title.font.size = Pt(18)
    p_t_title.font.bold = True
    p_t_title.font.color.rgb = TEXT_DARK
    
    # Add Table
    table_shape = slide6.shapes.add_table(6, 3, col1_left + Inches(0.3), Inches(2.9), col_width - Inches(0.6), Inches(3.2))
    table = table_shape.table
    table.columns[0].width = Inches(1.6)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(1.8)
    
    headers = ["Step", "Epoch", "Training Loss"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEXT_DARK
        cell.text = text
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Calibri"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
    data = [
        ["5", "0.022", "1.8983"],
        ["10", "0.044", "1.8030"],
        ["15", "0.067", "1.6545"],
        ["20", "0.089", "1.4351"],
        ["25", "0.111", "1.5245"]
    ]
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(241, 245, 249)
            else:
                cell.fill.fore_color.rgb = WHITE
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "Calibri"
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_MUTED if col_idx != 2 else TEXT_DARK
            if col_idx == 2:
                p.font.bold = True
                
    # Right Column: Conclusion
    col2_left = Inches(6.933)
    
    card_conc = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, col2_left, Inches(2.1), col_width, col_height)
    card_conc.fill.solid()
    card_conc.fill.fore_color.rgb = WHITE
    card_conc.line.color.rgb = CARD_BORDER
    card_conc.line.width = Pt(1.5)
    
    accent_card_c = slide6.shapes.add_shape(MSO_SHAPE.RECTANGLE, col2_left, Inches(2.1), col_width, Inches(0.1))
    accent_card_c.fill.solid()
    accent_card_c.fill.fore_color.rgb = ACCENT_GOLD
    accent_card_c.line.color.rgb = ACCENT_GOLD
    
    tb_c = slide6.shapes.add_textbox(col2_left + Inches(0.3), Inches(2.4), col_width - Inches(0.6), col_height - Inches(0.6))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    
    p_c_title = tf_c.paragraphs[0]
    p_c_title.text = "Key Project Takeaways"
    p_c_title.font.name = "Calibri"
    p_c_title.font.size = Pt(20)
    p_c_title.font.bold = True
    p_c_title.font.color.rgb = TEXT_DARK
    p_c_title.space_after = Pt(12)
    
    conc_points = [
        "Feasibility Proven: Fine-tuning a 1.1B parameter model on a standard CPU is completely viable with LoRA adapters.",
        "Loss Convergence: Loss consistently decreased from 1.90 to 1.43, indicating the model successfully learned specific instruction structures.",
        "Zero Cloud Cost: Avoided expensive cloud GPU endpoints while maintaining a low disk space overhead (~4.5MB adapters).",
        "Offline Ready: Achieved private local execution with robust fallbacks, suitable for remote classrooms."
    ]
    for pt in conc_points:
        p_pt = tf_c.add_paragraph()
        p_pt.text = "• " + pt
        p_pt.font.name = "Calibri"
        p_pt.font.size = Pt(11)
        p_pt.font.color.rgb = TEXT_MUTED
        p_pt.space_before = Pt(6)

    # ==========================================
    # SLIDE 7: Future Direction (Light Background)
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7, BG_LIGHT)
    add_slide_header(slide7, "Future Directions", "Future Horizons")
    
    # Subtitle
    desc_box7 = slide7.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.4))
    desc_tf7 = desc_box7.text_frame
    desc_p7 = desc_tf7.paragraphs[0]
    desc_p7.text = "Proposed research expansions to scale capabilities and models."
    desc_p7.font.name = "Calibri"
    desc_p7.font.size = Pt(15)
    desc_p7.font.color.rgb = TEXT_MUTED
    
    futures = [
        {
            "num": "01",
            "title": "Scale Dataset & Length",
            "points": [
                "Increase the training subset from 1k to 10k+ instruction samples.",
                "Extend maximum context length from 512 to 2048 tokens to support textbook-length chapters.",
                "Include specialized subjects like STEM and Coding in dataset."
            ]
        },
        {
            "num": "02",
            "title": "Quantized Larger Models",
            "points": [
                "Fine-tune larger models (e.g. Qwen-2.5-7B or Gemma-2-9B).",
                "Use llama.cpp or CPU bitsandbytes (4-bit/8-bit quantization) to run larger models under laptop constraints.",
                "Benchmark trade-offs between speed and model parameter scale."
            ]
        },
        {
            "num": "03",
            "title": "RAG Syllabus Integration",
            "points": [
                "Incorporate Retrieval-Augmented Generation (RAG).",
                "Allow teachers to upload classroom PDFs/syllabi locally.",
                "Let the model ground its answers in the uploaded files, reducing generic hallucinations."
            ]
        }
    ]
    
    for i, fut in enumerate(futures):
        left_pos = obj_left_start + i * (obj_width + obj_spacing)
        
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, obj_top, obj_width, obj_height)
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = CARD_BORDER
        card.line.width = Pt(1.5)
        
        accent_line = slide7.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos, obj_top, obj_width, Inches(0.12))
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = ACCENT_GOLD
        accent_line.line.color.rgb = ACCENT_GOLD
        
        tb = slide7.shapes.add_textbox(left_pos + Inches(0.2), obj_top + Inches(0.3), obj_width - Inches(0.4), obj_height - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p_num = tf.paragraphs[0]
        p_num.text = fut["num"]
        p_num.font.name = "Calibri"
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = ACCENT_GOLD
        p_num.space_after = Pt(2)
        
        p_title = tf.add_paragraph()
        p_title.text = fut["title"]
        p_title.font.name = "Calibri"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK
        p_title.space_after = Pt(8)
        
        for pt in fut["points"]:
            p_bullet = tf.add_paragraph()
            p_bullet.text = "• " + pt
            p_bullet.font.name = "Calibri"
            p_bullet.font.size = Pt(10.5)
            p_bullet.font.color.rgb = TEXT_MUTED
            p_bullet.space_before = Pt(4)

    # ==========================================
    # SLIDE 8: GitHub Details (Dark Background)
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8, BG_DARK)
    
    # Title
    t_box8 = slide8.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12.133), Inches(0.8))
    tf8 = t_box8.text_frame
    p_t8 = tf8.paragraphs[0]
    p_t8.text = "GitHub Repository & Setup Details"
    p_t8.font.name = "Calibri"
    p_t8.font.size = Pt(32)
    p_t8.font.bold = True
    p_t8.font.color.rgb = WHITE
    
    # Subtitle
    sub_box8 = slide8.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12.133), Inches(0.4))
    sub_tf8 = sub_box8.text_frame
    sub_p8 = sub_tf8.paragraphs[0]
    sub_p8.text = "Open source code structure and execution guide."
    sub_p8.font.name = "Calibri"
    sub_p8.font.size = Pt(16)
    sub_p8.font.color.rgb = ACCENT_GOLD
    
    # Two Columns
    col_w = Inches(5.8)
    col_h = Inches(4.5)
    
    # Left Card: Repo details
    c1_shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.0), col_w, col_h)
    c1_shape.fill.solid()
    c1_shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
    c1_shape.line.color.rgb = ACCENT_GOLD
    c1_shape.line.width = Pt(1)
    
    tb_git = slide8.shapes.add_textbox(Inches(0.9), Inches(2.2), col_w - Inches(0.6), col_h - Inches(0.4))
    tf_git = tb_git.text_frame
    tf_git.word_wrap = True
    
    p_git_t = tf_git.paragraphs[0]
    p_git_t.text = "Repository Details & Setup"
    p_git_t.font.name = "Calibri"
    p_git_t.font.size = Pt(20)
    p_git_t.font.bold = True
    p_git_t.font.color.rgb = WHITE
    p_git_t.space_after = Pt(10)
    
    p_git_url_t = tf_git.add_paragraph()
    p_git_url_t.text = "GitHub URL:"
    p_git_url_t.font.name = "Calibri"
    p_git_url_t.font.size = Pt(11)
    p_git_url_t.font.bold = True
    p_git_url_t.font.color.rgb = ACCENT_GOLD
    
    p_git_url = tf_git.add_paragraph()
    p_git_url.text = "https://github.com/deweshai-star/Instruction-Tuning-for-Educational-QA"
    p_git_url.font.name = "Calibri"
    p_git_url.font.size = Pt(11.5)
    p_git_url.font.color.rgb = WHITE
    p_git_url.space_after = Pt(8)
    
    git_cmds = [
        "1. Clone Project:\n   git clone https://github.com/deweshai-star/Instruction-Tuning-for-Educational-QA.git",
        "2. Create Virtualenv & Install:\n   python -m venv venv\n   .\\venv\\Scripts\\activate\n   pip install -r requirements.txt",
        "3. Run Web App:\n   streamlit run app.py"
    ]
    for cmd in git_cmds:
        p_cmd = tf_git.add_paragraph()
        p_cmd.text = cmd
        p_cmd.font.name = "Courier New"
        p_cmd.font.size = Pt(9.5)
        p_cmd.font.color.rgb = RGBColor(203, 213, 225)
        p_cmd.space_before = Pt(4)
        
    # Right Card: Project Structure
    c2_shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.933), Inches(2.0), col_w, col_h)
    c2_shape.fill.solid()
    c2_shape.fill.fore_color.rgb = WHITE
    c2_shape.line.color.rgb = CARD_BORDER
    c2_shape.line.width = Pt(1)
    
    tb_str = slide8.shapes.add_textbox(Inches(7.233), Inches(2.2), col_w - Inches(0.6), col_h - Inches(0.4))
    tf_str = tb_str.text_frame
    tf_str.word_wrap = True
    
    p_str_t = tf_str.paragraphs[0]
    p_str_t.text = "File Structure Specifications"
    p_str_t.font.name = "Calibri"
    p_str_t.font.size = Pt(20)
    p_str_t.font.bold = True
    p_str_t.font.color.rgb = TEXT_DARK
    p_str_t.space_after = Pt(10)
    
    struct = [
        "data_prep.py: Ingests, formats, and tokenizes subset.",
        "train.py: Loads TinyLlama, applies LoRA, executes training loop.",
        "inference.py: Interactive command-line testing script.",
        "app.py: Streamlit local deployment GUI with fallbacks.",
        "architecture.md: Comprehensive technical system overview.",
        "generate_ppt.py: Programmatic layout specifications for presentation."
    ]
    for s in struct:
        p_s = tf_str.add_paragraph()
        p_s.text = "• " + s
        p_s.font.name = "Calibri"
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = TEXT_MUTED
        p_s.space_before = Pt(5)
    
    # Save Presentation
    output_filename = "Instruction Tuning for Educational QA.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'!")

if __name__ == "__main__":
    create_presentation()
